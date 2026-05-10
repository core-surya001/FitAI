
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xC8, 0xA2, 0x6E)  # Sand/Gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xED, 0xE8)
CHARCOAL = RGBColor(0x3A, 0x3A, 0x4A)
GREEN = RGBColor(0x5C, 0x8A, 0x6E)
DARK2 = RGBColor(0x16, 0x21, 0x3E)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox

def add_bullet_box(slide, items, left, top, width, height, font_size=14, color=WHITE, title=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        run.font.name = "Calibri"
    for i, item in enumerate(items):
        if title or i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  ▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

# ─── Slide 1: Title ───────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
add_rect(sl, 0, 0, 13.33, 7.5, DARK2)
add_rect(sl, 0, 0, 5.5, 7.5, DARK)
add_rect(sl, 5.5, 0, 0.05, 7.5, ACCENT)
# Right accent block
add_rect(sl, 8.5, 1.5, 4.5, 4.5, RGBColor(0x0F,0x3D,0x33))
add_text(sl, "FitAI", 0.6, 1.5, 5, 1.4, 64, bold=True, color=WHITE)
add_text(sl, "AI-Powered Virtual Try-On Platform", 0.6, 3.0, 5, 0.7, 20, color=ACCENT)
add_text(sl, "Transforming the way people shop for fashion\nthrough intelligent, real-time garment visualization.", 0.6, 3.9, 4.8, 1.2, 14, color=LIGHT_GRAY)
add_text(sl, "✦  Next.js  ·  FastAPI  ·  IDM-VTON  ·  Google Gemini", 0.6, 5.8, 5.5, 0.5, 12, color=ACCENT, italic=True)
add_text(sl, "Virtual Try-On  |  AI Stylist  |  Digital Wardrobe", 9.0, 3.1, 3.8, 0.5, 13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Fashion meets AI", 9.2, 3.7, 3.4, 0.5, 11, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

# ─── Slide 2: Problem Statement ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, LIGHT_GRAY)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "02  |  The Problem", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)
add_text(sl, "Why does virtual fashion shopping still fail consumers?", 0.5, 1.3, 12, 0.6, 17, color=CHARCOAL, italic=True)

cards = [
    ("🛒", "Returns Crisis", "~30% of online clothing purchases are returned, mostly due to fit or appearance mismatch."),
    ("📷", "Blind Buying", "Shoppers cannot visualize garments on their own body — they rely on models who look nothing like them."),
    ("⏱️", "Time Wasted", "Hours spent in fitting rooms or waiting for deliveries just to discover the item doesn't suit them."),
    ("💸", "Revenue Loss", "Retailers lose billions annually in reverse logistics and restocking costs from avoidable returns."),
]
for i, (icon, title, desc) in enumerate(cards):
    col = i * 3.1 + 0.5
    add_rect(sl, col, 2.1, 2.9, 4.2, DARK)
    add_text(sl, icon, col + 0.15, 2.3, 2.6, 0.8, 28)
    add_text(sl, title, col + 0.15, 3.2, 2.6, 0.5, 14, bold=True, color=ACCENT)
    add_text(sl, desc, col + 0.15, 3.8, 2.6, 2.2, 11.5, color=LIGHT_GRAY)

# ─── Slide 3: Solution ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK2)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "03  |  Our Solution — FitAI", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)
add_rect(sl, 0.5, 1.3, 12.3, 0.08, ACCENT)
add_text(sl, "FitAI bridges the gap between online browsing and real-world fitting using cutting-edge AI.", 0.5, 1.55, 12, 0.6, 15, color=LIGHT_GRAY, italic=True)

steps = [
    ("01", "Upload", "User uploads a full-body or half-body photo of themselves."),
    ("02", "Select Garment", "User picks any garment from the wardrobe or catalog."),
    ("03", "AI Generation", "IDM-VTON model realistically superimposes the garment onto the user's photo."),
    ("04", "Style & Save", "View the result, ask the AI Stylist for advice, and save to wardrobe."),
]
for i, (num, step, desc) in enumerate(steps):
    x = i * 3.1 + 0.5
    add_rect(sl, x, 2.4, 2.9, 4.0, RGBColor(0x0F,0x2A,0x3D))
    add_text(sl, num, x + 0.15, 2.55, 1, 0.6, 28, bold=True, color=ACCENT)
    add_text(sl, step, x + 0.15, 3.2, 2.6, 0.5, 15, bold=True, color=WHITE)
    add_text(sl, desc, x + 0.15, 3.75, 2.6, 2.4, 12, color=LIGHT_GRAY)

# ─── Slide 4: Key Features ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, LIGHT_GRAY)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "04  |  Key Features", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

feats = [
    ("🪄", "Virtual Try-On Engine", "Powered by IDM-VTON (Hugging Face). Realistic garment diffusion adapted to user pose and body shape."),
    ("🤖", "AI Stylist — Gemini", "Integrated chat assistant powered by Google Gemini for outfit recommendations and fashion tips."),
    ("👗", "Digital Wardrobe", "Users upload and manage their photo library and garment catalog in a personal digital closet."),
    ("💳", "Credit Economy", "Free, Pro, and Unlimited subscription tiers with a mock payment gateway for upgrades."),
    ("🖼️", "Lookbook Gallery", "Save, view, download, and delete previously generated AI looks in a beautiful gallery."),
    ("🌐", "Modern Web App", "Glassmorphism UI built with Next.js 14, Tailwind CSS, and Framer Motion animations."),
]
for i, (icon, title, desc) in enumerate(feats):
    row, col = divmod(i, 3)
    x = col * 4.2 + 0.5
    y = row * 2.5 + 1.4
    add_rect(sl, x, y, 4.0, 2.2, DARK)
    add_text(sl, icon + "  " + title, x + 0.2, y + 0.2, 3.6, 0.5, 13, bold=True, color=ACCENT)
    add_text(sl, desc, x + 0.2, y + 0.75, 3.6, 1.3, 11.5, color=LIGHT_GRAY)

# ─── Slide 5: Technology Stack ───────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
add_rect(sl, 0, 0, 13.33, 1.1, DARK2)
add_text(sl, "05  |  Technology Stack", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

cols = [
    ("Frontend", ["Next.js 14 (App Router)", "TypeScript", "Tailwind CSS", "Framer Motion", "Zustand (State)", "Axios", "Lucide Icons"]),
    ("Backend", ["FastAPI (Python)", "SQLAlchemy ORM", "JWT Authentication", "Bcrypt Password Hash", "SQLite / PostgreSQL", "Uvicorn / Gunicorn"]),
    ("AI / ML", ["IDM-VTON (HuggingFace)", "Google Gemini API", "Gradio Client", "Replicate API (MVP)", "OpenPose (Preprocessing)", "Rembg (BG Removal)"]),
    ("Infra & Cloud", ["Vercel (Frontend)", "Render.com (Backend)", "AWS S3 / Cloudflare R2", "PostgreSQL (Supabase)", "Redis + Celery (Queue)", "RunPod / EC2 GPU"]),
]
for i, (title, items) in enumerate(cols):
    x = i * 3.1 + 0.5
    add_rect(sl, x, 1.3, 3.0, 5.8, RGBColor(0x0D,0x1B,0x2A))
    add_rect(sl, x, 1.3, 3.0, 0.5, ACCENT)
    add_text(sl, title, x + 0.15, 1.33, 2.7, 0.45, 14, bold=True, color=DARK)
    add_bullet_box(sl, items, x + 0.1, 1.95, 2.8, 4.8, font_size=11.5, color=LIGHT_GRAY)

# ─── Slide 6: System Architecture ───────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK2)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "06  |  System Architecture", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

# Draw architecture nodes as boxes with arrows
nodes = [
    (0.4, 3.0, 2.5, 0.9, "🖥️  Frontend\nNext.js / React", RGBColor(0x1B,0x4F,0x72)),
    (3.5, 3.0, 2.5, 0.9, "⚡  Backend\nFastAPI (Python)", RGBColor(0x17,0x5E,0x42)),
    (6.6, 1.5, 2.5, 0.9, "🗄️  Database\nPostgreSQL", RGBColor(0x4A,0x23,0x5A)),
    (6.6, 3.0, 2.5, 0.9, "☁️  Cloud Storage\nAWS S3 / R2", RGBColor(0x7D,0x4A,0x10)),
    (6.6, 4.5, 2.5, 0.9, "📬  Task Queue\nRedis + Celery", RGBColor(0x1A,0x44,0x5C)),
    (9.7, 3.0, 2.5, 0.9, "🤖  ML Worker\nIDM-VTON GPU", RGBColor(0x3B,0x1A,0x57)),
    (9.7, 1.5, 2.5, 0.9, "🔗  HuggingFace\nModel Hub", RGBColor(0x2C,0x47,0x1E)),
]
for (x, y, w, h, label, clr) in nodes:
    add_rect(sl, x, y, w, h, clr)
    add_text(sl, label, x + 0.1, y + 0.05, w - 0.2, h - 0.1, 11, bold=False, color=WHITE)

# Arrow labels
arrows = [
    (3.0, 3.4, "REST API →"),
    (6.2, 1.9, "↑ CRUD"),
    (6.2, 3.4, "↑ Upload"),
    (6.2, 4.9, "↑ Enqueue"),
    (9.3, 3.4, "← Process"),
    (9.7, 2.0, "↑ Fetch"),
]
for (ax, ay, alabel) in arrows:
    add_text(sl, alabel, ax, ay, 1.5, 0.4, 10, color=ACCENT, italic=True)

add_text(sl, "Decoupled microservice architecture ensures scalability: heavy AI tasks are offloaded to GPU workers via Redis, keeping the API responsive.", 0.5, 6.55, 12.3, 0.7, 11.5, color=LIGHT_GRAY, italic=True)

# ─── Slide 7: Database Schema ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, LIGHT_GRAY)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "07  |  Database Schema", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

tables = [
    ("Users", ["id  (PK)", "email", "password_hash", "credits", "subscription_tier", "created_at"]),
    ("UserPhotos", ["id  (PK)", "user_id  (FK)", "s3_url", "label", "created_at"]),
    ("Garments", ["id  (PK)", "user_id  (FK, optional)", "s3_url", "category", "name", "created_at"]),
    ("TryOnJobs", ["id  (PK)", "user_id  (FK)", "user_photo_id  (FK)", "garment_id  (FK)", "status", "result_image_url", "created_at"]),
]
for i, (tname, fields) in enumerate(tables):
    x = i * 3.1 + 0.5
    add_rect(sl, x, 1.3, 2.9, 5.7, DARK)
    add_rect(sl, x, 1.3, 2.9, 0.5, GREEN)
    add_text(sl, tname, x + 0.15, 1.33, 2.6, 0.45, 13, bold=True, color=WHITE)
    add_bullet_box(sl, fields, x + 0.1, 1.95, 2.7, 4.8, font_size=11.5, color=LIGHT_GRAY)

add_text(sl, "ORM: SQLAlchemy  |  Dev: SQLite  |  Prod: PostgreSQL (Supabase / AWS RDS)", 0.5, 7.0, 12.3, 0.4, 11, color=CHARCOAL, italic=True, align=PP_ALIGN.CENTER)

# ─── Slide 8: AI/ML Pipeline ─────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK)
add_rect(sl, 0, 0, 13.33, 1.1, DARK2)
add_text(sl, "08  |  AI / ML Pipeline", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

stages = [
    ("Pre-processing", ["Background removal\n(rembg)", "Pose estimation\n(OpenPose)", "Human parsing\n(Segmentation)", "Cloth mask\ngeneration"]),
    ("Inference\n(IDM-VTON)", ["User image +\ngarment image", "Pose map +\nhuman parse input", "Diffusion model\ngarment transfer", "Adapts to pose,\nlighting & shape"]),
    ("Post-processing", ["HD upscaling\n(optional)", "Result stored\nin AWS S3", "Job status\nupdated in DB", "Frontend polls\nresult via API"]),
]
for i, (stage, steps) in enumerate(stages):
    x = i * 4.2 + 0.5
    add_rect(sl, x, 1.3, 3.9, 5.7, RGBColor(0x0D,0x1B,0x2A))
    add_rect(sl, x, 1.3, 3.9, 0.55, ACCENT)
    add_text(sl, stage, x + 0.15, 1.33, 3.6, 0.52, 14, bold=True, color=DARK)
    for j, step in enumerate(steps):
        sy = 2.1 + j * 1.2
        add_rect(sl, x + 0.2, sy, 3.5, 1.0, RGBColor(0x16,0x2D,0x42))
        add_text(sl, step, x + 0.35, sy + 0.08, 3.2, 0.85, 11.5, color=WHITE)

add_text(sl, "MVP Shortcut: Use Replicate.com API to access pre-hosted IDM-VTON — no GPU setup needed for initial launch.", 0.5, 7.0, 12.3, 0.4, 11.5, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)

# ─── Slide 9: Development Roadmap ────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, LIGHT_GRAY)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "09  |  Development Roadmap", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

phases = [
    ("Phase 1\nWeeks 1–2", "Foundation & UI", ["Initialize Next.js & FastAPI repos", "PostgreSQL setup + User auth (JWT)", "Core UI pages: Login, Dashboard, Closet"]),
    ("Phase 2\nWeeks 3–4", "Storage & APIs", ["AWS S3 integration for image uploads", "Garment & Photo CRUD API endpoints", "Connect frontend to backend APIs"]),
    ("Phase 3\nWeeks 5–6", "AI/ML Integration", ["Replicate / HuggingFace Try-On API", "Celery+Redis background task queue", "Frontend polling / result display"]),
    ("Phase 4\nWeeks 7–8", "Polish & Deploy", ["Gallery view for past try-ons", "Error handling & edge cases", "Deploy: Vercel + Render + Supabase"]),
]
for i, (phase, title, items) in enumerate(phases):
    x = i * 3.1 + 0.5
    add_rect(sl, x, 1.3, 3.0, 0.85, DARK)
    add_text(sl, phase, x + 0.15, 1.33, 2.7, 0.82, 11.5, bold=True, color=ACCENT)
    add_rect(sl, x, 2.2, 3.0, 4.7, RGBColor(0x2A,0x2A,0x3E))
    add_text(sl, title, x + 0.15, 2.3, 2.7, 0.5, 14, bold=True, color=WHITE)
    add_bullet_box(sl, items, x + 0.1, 2.9, 2.8, 3.8, font_size=12, color=LIGHT_GRAY)

add_text(sl, "Total Duration: ~8 Weeks  |  Agile Sprint-based development", 0.5, 7.05, 12.3, 0.35, 11, color=CHARCOAL, italic=True, align=PP_ALIGN.CENTER)

# ─── Slide 10: Summary & Future ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, DARK2)
add_rect(sl, 0, 0, 13.33, 1.1, DARK)
add_text(sl, "10  |  Summary & Future Scope", 0.5, 0.18, 10, 0.7, 26, bold=True, color=WHITE)

add_rect(sl, 0.5, 1.25, 6.0, 5.7, RGBColor(0x0D,0x1B,0x2A))
add_text(sl, "🎯  What We've Built", 0.7, 1.35, 5.6, 0.5, 15, bold=True, color=ACCENT)
summary_items = [
    "Full-stack Virtual Try-On web platform",
    "IDM-VTON AI model integration via HuggingFace",
    "Google Gemini-powered AI Stylist chatbot",
    "JWT-secured REST API with FastAPI",
    "Digital Wardrobe with S3 image management",
    "Credit & subscription economy",
    "Production deployments: Vercel + Render",
    "Glassmorphism UI with animations (Next.js 14)",
]
add_bullet_box(sl, summary_items, 0.7, 1.95, 5.6, 4.7, font_size=12.5, color=LIGHT_GRAY)

add_rect(sl, 7.0, 1.25, 5.8, 5.7, RGBColor(0x0F,0x2E,0x1F))
add_text(sl, "🚀  Future Scope", 7.2, 1.35, 5.4, 0.5, 15, bold=True, color=ACCENT)
future_items = [
    "Mobile app (React Native / Flutter)",
    "3D body avatar generation",
    "Real-time AR try-on via webcam",
    "Social sharing & community lookbook",
    "Brand/retailer API integration for live catalog",
    "Multi-garment outfit layering",
    "Size recommendation engine (ML-based)",
    "NFT-based digital wardrobe ownership",
]
add_bullet_box(sl, future_items, 7.2, 1.95, 5.4, 4.7, font_size=12.5, color=LIGHT_GRAY)

add_rect(sl, 0, 7.05, 13.33, 0.45, ACCENT)
add_text(sl, "FitAI  —  The Future of Fashion is Virtual", 0.5, 7.08, 12.3, 0.38, 13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Save
out_path = r"c:\Users\Surya\Downloads\New folder\FitAI_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
